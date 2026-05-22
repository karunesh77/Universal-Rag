# =====================================================
# FILE PROCESSOR - Text Extraction & Chunking
# =====================================================
# Ye file different document types se text nikalti hai
# PDF, Word, Excel, PowerPoint, TXT - sab handle karta hai
#
# Text kyu nikalte hain?
# ─────────────────────────
# RAG System mein:
# 1. User document upload karta hai (PDF, Word, etc)
# 2. Hum text nikalte hain document se
# 3. Text ko chhote pieces (chunks) mein divide karte hain
# 4. Har chunk ka embedding (vector) banate hain
# 5. Jab user question pooche → closest chunks dhundho
# 6. Claude ko context deke answer generate karo
# =====================================================

import os
import logging
from typing import List, Optional

# Logging setup (errors track karne ke liye)
logger = logging.getLogger(__name__)


# =====================================================
# SECTION 1: PDF TEXT EXTRACTION
# =====================================================

def extract_text_from_pdf(file_path: str) -> str:
    """
    PDF file se text nikalna

    Strategy:
    1. Pehle pdfplumber try karo (better accuracy)
    2. Agar fail ho to PyPDF2 use karo (fallback)

    Params:
        file_path: PDF file ki local path

    Returns:
        Extracted text as string (empty string if failed)
    """

    text = ""

    # Strategy 1: pdfplumber (more accurate, handles complex PDFs)
    try:
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            pages_text = []
            for page_num, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    pages_text.append(f"[Page {page_num + 1}]\n{page_text}")

            text = "\n\n".join(pages_text)

        if text.strip():
            logger.info(f"PDF extracted with pdfplumber: {len(text)} chars")
            return text.strip()

    except ImportError:
        logger.warning("pdfplumber not installed, trying PyPDF2...")
    except Exception as e:
        logger.warning(f"pdfplumber failed: {e}, trying PyPDF2...")

    # Strategy 2: PyPDF2 (fallback)
    try:
        import PyPDF2
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            pages_text = []
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    pages_text.append(f"[Page {page_num + 1}]\n{page_text}")

            text = "\n\n".join(pages_text)

        logger.info(f"PDF extracted with PyPDF2: {len(text)} chars")
        return text.strip()

    except ImportError:
        logger.error("Neither pdfplumber nor PyPDF2 installed!")
        raise RuntimeError("PDF processing library not available. Run: pip install pdfplumber")
    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        raise RuntimeError(f"Could not extract text from PDF: {str(e)}")


# =====================================================
# SECTION 2: WORD DOCUMENT TEXT EXTRACTION
# =====================================================

def extract_text_from_docx(file_path: str) -> str:
    """
    Word document (.docx) se text nikalna

    Word documents mein:
    - Paragraphs: Normal text
    - Tables: Data in rows/columns
    - Headers/Footers: Page ke top/bottom text

    Params:
        file_path: .docx file ki local path

    Returns:
        Extracted text as string
    """

    try:
        from docx import Document as DocxDocument

        doc = DocxDocument(file_path)
        text_parts = []

        # Paragraphs se text nikalo
        for para in doc.paragraphs:
            if para.text.strip():
                # Heading hai to marker add karo
                if para.style.name.startswith('Heading'):
                    text_parts.append(f"\n## {para.text.strip()}")
                else:
                    text_parts.append(para.text.strip())

        # Tables se text nikalo
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_cells:
                    table_rows.append(" | ".join(row_cells))

            if table_rows:
                text_parts.append("\n[Table]\n" + "\n".join(table_rows))

        result = "\n".join(text_parts)
        logger.info(f"DOCX extracted: {len(result)} chars")
        return result.strip()

    except ImportError:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        raise RuntimeError(f"Could not extract text from Word document: {str(e)}")


# =====================================================
# SECTION 3: EXCEL SPREADSHEET TEXT EXTRACTION
# =====================================================

def extract_text_from_xlsx(file_path: str) -> str:
    """
    Excel spreadsheet (.xlsx) se text nikalna

    Excel mein:
    - Multiple sheets hoti hain
    - Rows aur columns mein data hota hai
    - Cells mein text, numbers, formulas hote hain

    Params:
        file_path: .xlsx file ki local path

    Returns:
        Extracted text as string (sheet name + data)
    """

    try:
        import openpyxl

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"\n[Sheet: {sheet_name}]")

            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                # Row cells ko string mein convert karo
                row_cells = []
                for cell in row:
                    if cell is not None and str(cell).strip():
                        row_cells.append(str(cell))
                    else:
                        row_cells.append("")

                # Empty rows skip karo
                if any(c.strip() for c in row_cells):
                    text_parts.append("\t".join(row_cells))
                    row_count += 1

            logger.info(f"Sheet '{sheet_name}': {row_count} rows extracted")

        wb.close()
        result = "\n".join(text_parts)
        logger.info(f"XLSX extracted: {len(result)} chars")
        return result.strip()

    except ImportError:
        raise RuntimeError("openpyxl not installed. Run: pip install openpyxl")
    except Exception as e:
        logger.error(f"XLSX extraction failed: {e}")
        raise RuntimeError(f"Could not extract text from Excel: {str(e)}")


# =====================================================
# SECTION 4: POWERPOINT TEXT EXTRACTION
# =====================================================

def extract_text_from_pptx(file_path: str) -> str:
    """
    PowerPoint presentation (.pptx) se text nikalna

    PowerPoint mein:
    - Multiple slides hoti hain
    - Har slide mein shapes hoti hain
    - Shapes mein text boxes, titles, bullets hote hain

    Params:
        file_path: .pptx file ki local path

    Returns:
        Extracted text as string (slide number + content)
    """

    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            text_parts.append(f"\n[Slide {slide_num}]")

            for shape in slide.shapes:
                # Shape mein text hai?
                if hasattr(shape, "text") and shape.text.strip():
                    # Title vs normal text
                    if shape.shape_type == 13:  # Title shape
                        slide_texts.append(f"Title: {shape.text.strip()}")
                    else:
                        slide_texts.append(shape.text.strip())

                # Table shape hai?
                if shape.has_table:
                    for row in shape.table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_cells:
                            slide_texts.append(" | ".join(row_cells))

            if slide_texts:
                text_parts.extend(slide_texts)

        result = "\n".join(text_parts)
        logger.info(f"PPTX extracted: {len(result)} chars ({len(prs.slides)} slides)")
        return result.strip()

    except ImportError:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        raise RuntimeError(f"Could not extract text from PowerPoint: {str(e)}")


# =====================================================
# SECTION 5: PLAIN TEXT FILE READING
# =====================================================

def extract_text_from_txt(file_path: str) -> str:
    """
    Plain text file (.txt) padna

    TXT files mein directly text hota hai
    Bas read karke return karo

    Params:
        file_path: .txt file ki local path

    Returns:
        File content as string
    """

    try:
        # Try UTF-8 pehle (most common)
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        logger.info(f"TXT extracted: {len(content)} chars")
        return content.strip()

    except UnicodeDecodeError:
        # UTF-8 fail hua, try with error handling
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            logger.warning(f"TXT extracted with UTF-8 errors ignored: {len(content)} chars")
            return content.strip()
        except Exception as e:
            raise RuntimeError(f"Could not read text file: {str(e)}")

    except Exception as e:
        logger.error(f"TXT extraction failed: {e}")
        raise RuntimeError(f"Could not read text file: {str(e)}")


# =====================================================
# SECTION 6: MAIN EXTRACT FUNCTION
# =====================================================

def extract_text(file_path: str, file_type: str) -> str:
    """
    File type ke basis par sahi extractor call karo

    Ye function ek dispatcher hai:
    - file_type "pdf" → extract_text_from_pdf()
    - file_type "docx" → extract_text_from_docx()
    - file_type "xlsx" → extract_text_from_xlsx()
    - file_type "pptx" → extract_text_from_pptx()
    - file_type "txt" → extract_text_from_txt()

    Params:
        file_path: File ki path (server par)
        file_type: File ka type (lowercase string)

    Returns:
        Extracted text

    Raises:
        ValueError: Agar unsupported file type diya
        RuntimeError: Agar extraction fail ho
    """

    # Available extractors ka map
    extractors = {
        'pdf':  extract_text_from_pdf,
        'docx': extract_text_from_docx,
        'xlsx': extract_text_from_xlsx,
        'pptx': extract_text_from_pptx,
        'txt':  extract_text_from_txt,
    }

    file_type_lower = file_type.lower().strip()

    if file_type_lower not in extractors:
        supported = ", ".join(extractors.keys())
        raise ValueError(f"Unsupported file type '{file_type}'. Supported: {supported}")

    # Sahi extractor call karo
    extractor_fn = extractors[file_type_lower]
    return extractor_fn(file_path)


# =====================================================
# SECTION 7: TEXT CHUNKING
# =====================================================

def chunk_text(
    text: str,
    chunk_size: int = 1000,
    overlap: int = 200
) -> List[str]:
    """
    Text ko chhote pieces (chunks) mein divide karna

    Kyu chunking zaroori hai?
    ─────────────────────────
    - AI models ka context window limited hota hai
    - Ek baar mein poora document nahi bhej sakte
    - Relevant parts dhundho → sirf woh bhejo

    Algorithm:
    ──────────
    text = "ABCDE...XYZ" (1000 chars)
    chunk_size = 500, overlap = 100

    Chunk 1: chars 0-500 → "ABCDE...E"
    Chunk 2: chars 400-900 → "BCDE...Y"  (100 overlap start)
    Chunk 3: chars 800-1000 → "XYZ"

    Overlap kyu?
    ─────────────
    Agar ek sentence chunk boundary par ho, wo dono chunks mein hoga
    Matlab context nahi katega

    Params:
        text: Input text (poora document)
        chunk_size: Har chunk ki max size (characters)
        overlap: Adjacent chunks ke beech shared text

    Returns:
        List of text chunks
    """

    if not text or not text.strip():
        return []

    text = text.strip()
    chunks = []
    start = 0
    text_length = len(text)

    while start < text_length:
        # Chunk ka end position
        end = min(start + chunk_size, text_length)
        chunk = text[start:end]

        # Natural break point dhundho (sentence ya paragraph end)
        if end < text_length:
            # Paragraph break prefer karo
            paragraph_break = chunk.rfind('\n\n')
            sentence_break = chunk.rfind('. ')
            newline_break = chunk.rfind('\n')

            # Best break point (prefer longer break, but at least half chunk)
            min_break = chunk_size // 2

            if paragraph_break > min_break:
                chunk = chunk[:paragraph_break + 2]
                end = start + paragraph_break + 2
            elif sentence_break > min_break:
                chunk = chunk[:sentence_break + 2]
                end = start + sentence_break + 2
            elif newline_break > min_break:
                chunk = chunk[:newline_break + 1]
                end = start + newline_break + 1

        chunk = chunk.strip()
        if chunk:
            chunks.append(chunk)

        # Next chunk start (overlap ke saath)
        start = end - overlap

        # Prevent infinite loop agar progress nahi ho raha
        if start >= end:
            start = end

    # Empty chunks remove karo
    chunks = [c for c in chunks if c.strip()]

    logger.info(f"Text chunked: {text_length} chars → {len(chunks)} chunks "
                f"(size={chunk_size}, overlap={overlap})")
    return chunks


# =====================================================
# SECTION 8: FILE VALIDATION
# =====================================================

# Allowed file types
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'xlsx', 'pptx', 'txt'}

# Max file size: 50MB (bytes mein)
MAX_FILE_SIZE = 50 * 1024 * 1024  # 52,428,800 bytes


def validate_file(filename: str, file_size: int) -> tuple[bool, str]:
    """
    File validate karna upload se pehle

    Params:
        filename: File ka naam (extension ke saath)
        file_size: File ki size in bytes

    Returns:
        (is_valid: bool, error_message: str)
        is_valid = True matlab file theek hai
        error_message = kya galat hai (agar kuch galat ho)
    """

    # Extension nikalo filename se
    if '.' not in filename:
        return False, "File extension nahi hai (e.g., document.pdf)"

    extension = filename.rsplit('.', 1)[1].lower()

    # Extension check
    if extension not in ALLOWED_EXTENSIONS:
        allowed = ", ".join(ALLOWED_EXTENSIONS)
        return False, f"File type '{extension}' allowed nahi. Allowed: {allowed}"

    # File size check
    if file_size <= 0:
        return False, "File empty hai"

    if file_size > MAX_FILE_SIZE:
        max_mb = MAX_FILE_SIZE / (1024 * 1024)
        actual_mb = file_size / (1024 * 1024)
        return False, f"File too large ({actual_mb:.1f}MB). Max allowed: {max_mb:.0f}MB"

    return True, ""


def get_file_extension(filename: str) -> str:
    """
    Filename se extension nikalna

    Examples:
        "report.pdf" → "pdf"
        "data.xlsx" → "xlsx"
        "notes.txt" → "txt"
    """
    if '.' not in filename:
        raise ValueError(f"No extension in filename: {filename}")
    return filename.rsplit('.', 1)[1].lower()


def format_file_size(size_bytes: int) -> str:
    """
    Bytes ko human-readable format mein convert karna

    Examples:
        1024 → "1.0 KB"
        1048576 → "1.0 MB"
        5242880 → "5.0 MB"
    """
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    else:
        return f"{size_bytes / (1024 * 1024):.1f} MB"
